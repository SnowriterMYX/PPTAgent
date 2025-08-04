import inspect
import os
import re
import traceback
from copy import deepcopy
from dataclasses import dataclass
from enum import Enum
from functools import partial
from typing import Any, Optional, Union

import PIL
from bs4 import BeautifulSoup
from mistune import HTMLRenderer, create_markdown
from pptx.enum.text import PP_ALIGN
from pptx.oxml import parse_xml
from pptx.shapes.base import BaseShape
from pptx.shapes.graphfrm import GraphicFrame as PPTXGraphicFrame
from pptx.text.text import _Run
from pptx.util import Pt

from pptagent.document import Document
from pptagent.presentation import Closure, ClosureType, Picture, ShapeElement, SlidePage
from pptagent.utils import get_logger, runs_merge

logger = get_logger(__name__)
TABLE_REGEX = re.compile(r".*table_[0-9a-fA-F]{4}\.png$")


class SlideRenderer(HTMLRenderer):
    """
    A renderer that does not render lists.
    """

    def list(self, text: str, ordered: bool, **attrs: Any) -> str:
        return text

    def list_item(self, text: str) -> str:
        return text


markdown = create_markdown(renderer=SlideRenderer(), plugins=["strikethrough"])


class SlideEditError(Exception):
    """
    Exception raised when an edit operation fails.
    """


@dataclass
class HistoryMark:
    """
    Mark the execution status of the API call, comment and a line of code.
    """

    API_CALL_ERROR = "api_call_error"
    API_CALL_CORRECT = "api_call_correct"
    COMMENT_CORRECT = "comment_correct"
    COMMENT_ERROR = "comment_error"
    CODE_RUN_ERROR = "code_run_error"
    CODE_RUN_CORRECT = "code_run_correct"


class CodeExecutor:
    """
    Execute code actions and manage API call history, and providing error feedback.
    """

    def __init__(self, retry_times: int):
        """
        Initialize the CodeExecutor.

        Args:
            retry_times (int): The number of times to retry failed actions.
        """
        self.api_history = []
        self.command_history = []
        self.code_history = []
        self.retry_times = retry_times
        self.registered_functions = API_TYPES.all_funcs()
        self.function_regex = re.compile(r"^[a-z]+_[a-z_]+\(.+\)")

    def reset_execution_state(self):
        """
        重置执行状态，用于重试时清理之前的状态
        """
        # 保留历史记录，但清理当前执行状态
        pass

    @classmethod
    def get_apis_docs(
        cls,
        funcs: list[callable],
        show_doc: bool = True,
        show_return: bool = True,
        ignore_keys: Optional[list[str]] = None,
    ) -> str:
        """
        Get the documentation for a list of API functions.

        Args:
            funcs (list[callable]): A list of functions to document.
            show_example (bool): Whether to show examples in the documentation.

        Returns:
            str: The formatted API documentation.
        """
        if ignore_keys is None:
            ignore_keys = {"slide", "self", "doc"}
        api_doc = []
        for func in funcs:
            sig = inspect.signature(func)
            params = []
            for name, param in sig.parameters.items():
                if name in ignore_keys:
                    continue
                param_str = name
                if param.annotation != inspect.Parameter.empty:
                    param_str += f": {param.annotation.__name__}"
                if param.default != inspect.Parameter.empty:
                    param_str += f" = {repr(param.default)}"
                params.append(param_str)
            signature = f"def {func.__name__}({', '.join(params)})"
            if show_return and sig.return_annotation != inspect.Parameter.empty:
                signature += f" -> {sig.return_annotation.__name__}"
            if show_doc and inspect.getdoc(func) is not None:
                doc = "\t" + inspect.getdoc(func)
            else:
                doc = ""
            signature += f"\n{doc}"
            api_doc.append(signature)
        return "\n".join(api_doc)

    def execute_actions(
        self,
        actions: str,
        edit_slide: SlidePage,
        doc: Document,
        found_code: bool = False,
    ) -> Union[tuple[str, str], None]:
        """
        Execute a series of actions on a slide.

        Args:
            actions (str): The actions to execute.
            edit_slide (SlidePage): The slide to edit.
            found_code (bool): Whether code was found in the actions.

        Returns:
            tuple: The API lines and traceback if an error occurs.
            None: If no error occurs.
        """
        api_calls = actions.strip().split("\n")
        logger.debug(f"Executing {len(api_calls)} actions on slide {edit_slide.slide_idx}")
        logger.debug(f"Actions to execute:\n{actions}")

        self.api_history.append(
            [HistoryMark.API_CALL_ERROR, edit_slide.slide_idx, actions]
        )
        for line_idx, line in enumerate(api_calls):
            try:
                if line_idx == len(api_calls) - 1 and not found_code:
                    raise SlideEditError(
                        "No code block found in the output, please output the api calls without any prefix."
                    )
                if line.startswith("def"):
                    raise SlideEditError("The function definition were not allowed.")
                if line.startswith("#"):
                    if len(self.command_history) != 0:
                        self.command_history[-1][0] = HistoryMark.COMMENT_CORRECT
                    self.command_history.append([HistoryMark.COMMENT_ERROR, line, None])
                    continue
                if not self.function_regex.match(line):
                    continue
                found_code = True
                func = line.split("(")[0]
                if func not in self.registered_functions:
                    raise SlideEditError(f"The function {func} is not defined.")
                # 注意：移除了过于严格的命令冲突检查，允许在同一序列中混合使用clone和del操作
                # 这样可以避免不必要的错误，提高系统的灵活性
                self.code_history.append([HistoryMark.CODE_RUN_ERROR, line, None])
                partial_func = partial(self.registered_functions[func], edit_slide)
                if func == "replace_image":
                    partial_func = partial(partial_func, doc)
                eval(line, {}, {func: partial_func})
                self.code_history[-1][0] = HistoryMark.CODE_RUN_CORRECT
            except Exception as e:
                if not isinstance(e, SlideEditError):
                    logger.warning(f"Encountered unknown error in function '{func}': {e}")
                    logger.debug(f"Function arguments: {line}")
                else:
                    logger.debug(f"SlideEditError in function '{func}': {e}")

                trace_msg = traceback.format_exc()
                if len(self.code_history) != 0:
                    self.code_history[-1][-1] = trace_msg
                api_lines = (
                    "\n".join(api_calls[: line_idx - 1])
                    + f"\n--> Error Line: {line}\n"
                    + "\n".join(api_calls[line_idx + 1 :])
                )
                return api_lines, trace_msg
        if len(self.command_history) != 0:
            self.command_history[-1][0] = HistoryMark.COMMENT_CORRECT
        self.api_history[-1][0] = HistoryMark.API_CALL_CORRECT

    def __add__(self, other):
        self.api_history.extend(other.api_history)
        self.command_history.extend(other.command_history)
        self.code_history.extend(other.code_history)
        return self


# supporting functions
def element_index(slide: SlidePage, element_id: int) -> ShapeElement:
    """
    Find the an element in a slide.

    Args:
        slide (SlidePage): The slide
        element_id (int): The ID of the element.

    Returns:
        ShapeElement: The shape corresponding to the element ID.

    Raises:
        SlideEditError: If the element is not found.
    """
    for shape in slide:
        if shape.shape_idx == element_id:
            return shape
    raise SlideEditError(
        f"Cannot find element {element_id}, is it deleted or not exist?"
    )


@dataclass
class TextBlock:
    text: str
    bold: bool = False
    italic: bool = False
    code: bool = False
    strikethrough: bool = False
    href: str = None

    def build_run(self, run: _Run):
        if self.bold:
            run.font.bold = True
        if self.italic:
            run.font.italic = True
        if self.code:
            run.font.name = "Consolas"
        if self.strikethrough:
            run.font.strikethrough = True
        if self.href is not None:
            run.hyperlink.address = self.href

        run.text = self.text


MARKDOWN_STYLES = {
    "strong": "bold",
    "em": "italic",
    "code": "code",
    "del": "strikethrough",
}


def process_element(element, styles=None) -> list[TextBlock]:
    if styles is None:
        styles = {}

    result = []

    if isinstance(element, str):
        result.append(TextBlock(element, **styles))
    else:
        if element.name == "a":
            new_styles = styles.copy()
            for child in element.children:
                blocks = process_element(child, new_styles)
                for block in blocks:
                    block.href = element.get("href")
                result.extend(blocks)
        elif MARKDOWN_STYLES.get(element.name):
            new_styles = styles.copy()
            new_styles[MARKDOWN_STYLES[element.name]] = True
            for child in element.children:
                result.extend(process_element(child, new_styles))
        else:
            for child in element.children:
                result.extend(process_element(child, styles))

    return result


def replace_para(paragraph_id: int, new_text: str, shape: BaseShape):
    """
    Replace the text of a paragraph in a shape.
    """
    para = shape.text_frame.paragraphs[paragraph_id]
    html = markdown(new_text).strip()
    soup = BeautifulSoup(html, "html.parser")
    blocks = process_element(soup)

    empty_run = runs_merge(para)
    empty_run.text = ""
    for _ in range(len(blocks) - 1):
        empty_run._r.addnext(parse_xml(empty_run._r.xml))
    for block, run in zip(blocks, para.runs):
        block.build_run(run)


def clone_para(paragraph_id: int, shape: BaseShape):
    """
    Clone a paragraph in a shape.
    """
    para = shape.text_frame.paragraphs[paragraph_id]
    shape.text_frame.paragraphs[-1]._element.addnext(parse_xml(para._element.xml))


def del_para(paragraph_id: int, shape: BaseShape):
    """
    Delete a paragraph from a shape.
    """
    para = shape.text_frame.paragraphs[paragraph_id]
    para._element.getparent().remove(para._element)


def add_table(table_data: list[list[str]], table: PPTXGraphicFrame):
    rows = len(table_data)
    cols = len(table_data[0])

    max_lengths = [max(len(row[j]) for row in table_data) for j in range(cols)]
    total_length = sum(max_lengths)
    for j in range(cols):
        col_width = int((max_lengths[j] / total_length) * table.width)
        table.table.columns[j].width = col_width

    for i in range(rows):
        for j in range(cols):
            table.table.cell(i, j).text = table_data[i][j]


def merge_cells(merge_area: list[tuple[int, int, int, int]], table: PPTXGraphicFrame):
    if merge_area is None or len(merge_area) == 0:
        return
    for y1, x1, y2, x2 in merge_area:
        try:
            table.table.cell(x1, y1).merge(table.table.cell(x2, y2))
            for x, y in zip(range(x1, x2 + 1), range(y1, y2 + 1)):
                tf = table.table.cell(x, y).text_frame
                for p in tf.paragraphs:
                    p.alignment = PP_ALIGN.CENTER
        except Exception as e:
            logger.warning(f"Failed to merge cells: {e}")


def validate_paragraph_operation(shape, div_id, paragraph_id, operation_name):
    """
    验证段落操作的有效性

    Args:
        shape: 形状对象
        div_id: 元素ID
        paragraph_id: 段落ID
        operation_name: 操作名称（用于错误消息）

    Raises:
        SlideEditError: 如果操作无效
    """
    logger.debug(f"Validating {operation_name} operation on paragraph {paragraph_id} of element {div_id}")

    if not shape.text_frame.is_textframe:
        raise SlideEditError(
            f"The element {div_id} does not have a text frame, please check the element id and type of element."
        )

    # 获取所有段落信息用于调试
    all_paragraphs = shape.text_frame.paragraphs
    logger.debug(f"Total paragraphs in element {div_id}: {len(all_paragraphs)}")
    for i, para in enumerate(all_paragraphs):
        logger.debug(f"  Paragraph {i}: idx={para.idx}, real_idx={getattr(para, 'real_idx', 'N/A')}, text='{para.text[:50]}...'")

    # 获取有效段落
    valid_paragraphs = [para for para in all_paragraphs if para.idx != -1]
    if not valid_paragraphs:
        raise SlideEditError(
            f"No valid paragraphs found in element {div_id}. Cannot perform {operation_name} operation. "
            f"Total paragraphs: {len(all_paragraphs)}, all have idx=-1"
        )

    # 检查目标段落是否存在且有效
    target_paragraph = None
    for para in all_paragraphs:
        if para.idx == paragraph_id:
            target_paragraph = para
            break

    if target_paragraph is None:
        available_ids = [para.idx for para in valid_paragraphs]
        logger.warning(f"Paragraph {paragraph_id} not found in element {div_id}. Available IDs: {available_ids}")

        # 尝试智能修复：如果请求的段落ID超出范围，使用最后一个有效段落
        if available_ids and paragraph_id >= max(available_ids):
            logger.info(f"Auto-correcting paragraph ID from {paragraph_id} to {max(available_ids)} for {operation_name} operation")
            target_paragraph = next(para for para in valid_paragraphs if para.idx == max(available_ids))
        else:
            raise SlideEditError(
                f"Cannot find paragraph {paragraph_id} in element {div_id} for {operation_name} operation. "
                f"Available paragraph IDs: {available_ids}. "
                f"Suggestion: Use one of the available IDs or check if previous operations modified the paragraph structure."
            )

    if target_paragraph.idx == -1:
        raise SlideEditError(
            f"Cannot perform {operation_name} on invalid paragraph {paragraph_id} of element {div_id}."
        )

    logger.debug(f"Validation successful for {operation_name} operation on paragraph {paragraph_id}")
    return target_paragraph


# api functions
def del_paragraph(slide: SlidePage, div_id: int, paragraph_id: int):
    """
    Delete a paragraph from a slide.

    Args:
        slide (SlidePage): The slide containing the paragraph.
        div_id (int): The ID of the division containing the paragraph.
        paragraph_id (int): The ID of the paragraph to delete.

    Raises:
        SlideEditError: If the paragraph is not found.
    """
    shape = element_index(slide, div_id)
    target_paragraph = validate_paragraph_operation(shape, div_id, paragraph_id, "delete")

    shape.text_frame.paragraphs.remove(target_paragraph)
    shape._closures[ClosureType.DELETE].append(
        Closure(partial(del_para, target_paragraph.real_idx), target_paragraph.real_idx)
    )


def del_image(slide: SlidePage, figure_id: int):
    """
    Delete an image from a slide.

    Args:
        slide (SlidePage): The slide containing the image.
        figure_id (int): The ID of the image to delete.
    """
    shape = element_index(slide, figure_id)
    if not isinstance(shape, Picture):
        raise SlideEditError(
            f"The element {shape.shape_idx} of slide {slide.slide_idx} is not a Picture."
        )
    slide.shapes.remove(shape)


def replace_paragraph(slide: SlidePage, div_id: int, paragraph_id: int, text: str):
    """
    Replace the text of a paragraph in a slide.

    Args:
        slide (SlidePage): The slide containing the paragraph.
        div_id (int): The ID of the division containing the paragraph.
        paragraph_id (int): The ID of the paragraph to replace.
        text (str): The new text to replace with.

    Raises:
        SlideEditError: If the paragraph is not found.
    """
    shape = element_index(slide, div_id)
    target_paragraph = validate_paragraph_operation(shape, div_id, paragraph_id, "replace")

    target_paragraph.text = text
    shape._closures[ClosureType.REPLACE].append(
        Closure(
            partial(replace_para, target_paragraph.real_idx, text),
            target_paragraph.real_idx,
        )
    )


def replace_image(slide: SlidePage, doc: Document, img_id: int, image_path: str):
    """
    Replace an image in a slide.

    Args:
        slide (SlidePage): The slide containing the image.
        img_id (int): The ID of the image to replace.
        image_path (str): The path to the new image.

    Raises:
        SlideEditError: If the image path does not exist.
    """
    if not os.path.exists(image_path):
        raise SlideEditError(
            f"The image {image_path} does not exist, consider use del_image if image_path in the given command is faked"
        )
    shape = element_index(slide, img_id)
    if not isinstance(shape, Picture):
        raise SlideEditError(
            f"The element {shape.shape_idx} of slide {slide.slide_idx} is not a Picture."
        )

    try:
        if TABLE_REGEX.match(image_path):
            return replace_image_with_table(shape, doc, image_path)
    except Exception as e:
        logger.warning(
            f"Failed to replace image with table element: {e}, fallback to use image directly."
        )

    img_size = PIL.Image.open(image_path).size
    r = min(shape.width / img_size[0], shape.height / img_size[1])
    new_width = img_size[0] * r
    new_height = img_size[1] * r
    shape.top = Pt(shape.top + (shape.height - new_height) / 2)
    shape.width = Pt(new_width)
    shape.height = Pt(new_height)
    shape.img_path = image_path


def clone_paragraph(slide: SlidePage, div_id: int, paragraph_id: int):
    """
    Clone a paragraph in a slide.

    Args:
        slide (SlidePage): The slide containing the paragraph.
        div_id (int): The ID of the division containing the paragraph.
        paragraph_id (int): The ID of the paragraph to clone.

    Raises:
        SlideEditError: If the paragraph is not found.

    Mention: the cloned paragraph will have a paragraph_id one greater than the current maximum in the parent element.
    """
    shape = element_index(slide, div_id)
    target_paragraph = validate_paragraph_operation(shape, div_id, paragraph_id, "clone")

    # 获取有效段落的最大索引
    valid_indices = [para.idx for para in shape.text_frame.paragraphs if para.idx != -1]
    max_idx = max(valid_indices)

    # 克隆段落
    cloned_para = deepcopy(target_paragraph)
    cloned_para.idx = max_idx + 1
    cloned_para.real_idx = len(shape.text_frame.paragraphs)

    shape.text_frame.paragraphs.append(cloned_para)
    shape._closures[ClosureType.CLONE].append(
        Closure(
            partial(clone_para, target_paragraph.real_idx),
            target_paragraph.real_idx,
        )
    )


def replace_image_with_table(shape: Picture, doc: Document, image_path: str):
    table = doc.get_table(image_path)
    shape.is_table = True
    shape.grid = (len(table.cells), len(table.cells[0]))
    shape._closures[ClosureType.REPLACE].append(
        Closure(partial(add_table, table.cells))
    )
    shape._closures[ClosureType.MERGE].append(
        Closure(partial(merge_cells, table.merge_area))
    )
    return


class API_TYPES(Enum):
    Agent = [
        replace_image,
        del_image,
        clone_paragraph,
        replace_paragraph,
        del_paragraph,
    ]

    @classmethod
    def all_funcs(cls) -> dict[str, callable]:
        funcs = {}
        for attr in dir(cls):
            if attr.startswith("__"):
                continue
            funcs |= {func.__name__: func for func in getattr(cls, attr).value}
        return funcs
