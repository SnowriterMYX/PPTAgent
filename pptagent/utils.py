import asyncio
import json
import logging
import os
import platform
import shutil
import subprocess
import tempfile
import traceback
from itertools import product
from shutil import which
from time import sleep, time
from typing import Any, Optional

import json_repair
import Levenshtein
from html2image import Html2Image
from mistune import html as markdown
try:
    from pdf2image import convert_from_path
    PDF2IMAGE_AVAILABLE = True
except ImportError:
    PDF2IMAGE_AVAILABLE = False
    logger.warning("pdf2image not available. Please install poppler-utils.")
from PIL import Image as PILImage
from pptx.dml.color import RGBColor
from pptx.oxml import parse_xml
from pptx.parts.image import Image
from pptx.shapes.group import GroupShape
from pptx.text.text import _Paragraph, _Run
from pptx.util import Length, Pt
from tenacity import RetryCallState, retry, stop_after_attempt, wait_fixed


def get_logger(name="pptagent", level=None):
    """
    Get a logger with the specified name and level.

    Args:
        name (str): The name of the logger.
        level (int): The logging level (default: logging.INFO).

    Returns:
        logging.Logger: A configured logger instance.
    """
    if level is None:
        # 处理日志级别配置，支持字符串和数字
        log_level = os.environ.get("LOG_LEVEL", "INFO")
        if isinstance(log_level, str):
            level_map = {
                "DEBUG": logging.DEBUG,
                "INFO": logging.INFO,
                "WARNING": logging.WARNING,
                "ERROR": logging.ERROR,
                "CRITICAL": logging.CRITICAL
            }
            level = level_map.get(log_level.upper(), logging.INFO)
        else:
            try:
                level = int(log_level)
            except (ValueError, TypeError):
                level = logging.INFO

    logger = logging.getLogger(name)
    logger.setLevel(level)

    # Check if the logger already has handlers to avoid duplicates
    if not logger.handlers:
        # Create console handler and set level
        console_handler = logging.StreamHandler()
        console_handler.setLevel(level)

        # Create formatter
        formatter = logging.Formatter(
            "%(asctime)s - %(levelname)s - %(filename)s:%(lineno)d - %(message)s"
        )
        console_handler.setFormatter(formatter)

        # Add handler to logger
        logger.addHandler(console_handler)

    return logger


logger = get_logger(__name__)

def get_soffice_command():
    """
    获取LibreOffice命令，兼容不同操作系统
    """
    # Windows下的可能路径
    if platform.system() == "Windows":
        possible_commands = [
            "soffice",  # 如果已添加到PATH
            "soffice.exe",
            r"C:\Program Files\LibreOffice\program\soffice.exe",
            r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
            # 便携版路径
            r".\LibreOffice\program\soffice.exe",
            r".\LibreOfficePortable\App\libreoffice\program\soffice.exe",
        ]

        for cmd in possible_commands:
            if which(cmd) or os.path.exists(cmd):
                return cmd
        return None
    else:
        # Linux/macOS
        return which("soffice")

# 检查LibreOffice是否可用
SOFFICE_CMD = get_soffice_command()
if SOFFICE_CMD is None:
    logging.warning("LibreOffice (soffice) is not found. PPT to images conversion will not work.")
    logging.warning("Please install LibreOffice or add it to your PATH.")
    if platform.system() == "Windows":
        logging.warning("Windows users can download LibreOffice from: https://www.libreoffice.org/download/download/")
else:
    logging.info(f"Found LibreOffice at: {SOFFICE_CMD}")

# Set of supported image extensions
IMAGE_EXTENSIONS: set[str] = {
    "bmp",
    "jpg",
    "jpeg",
    "pgm",
    "png",
    "ppm",
    "tif",
    "tiff",
    "webp",
}

# Common colors and measurements
BLACK = RGBColor(0, 0, 0)
YELLOW = RGBColor(255, 255, 0)
BLUE = RGBColor(0, 0, 255)
BORDER_LEN = Pt(2)
BORDER_OFFSET = Pt(2)
LABEL_LEN = Pt(24)
FONT_LEN = Pt(20)


def is_image_path(file: str) -> bool:
    """
    Check if a file path is an image based on its extension.

    Args:
        file (str): The file path to check.

    Returns:
        bool: True if the file is an image, False otherwise.
    """
    return file.split(".")[-1].lower() in IMAGE_EXTENSIONS


def runs_merge(paragraph: _Paragraph) -> Optional[_Run]:
    """
    Merge all runs in a paragraph into a single run.

    Args:
        paragraph (_Paragraph): The paragraph to merge runs in.

    Returns:
        Optional[_Run]: The merged run, or None if there are no runs.
    """
    runs = paragraph.runs

    # Handle field codes
    if len(runs) == 0:
        runs = [
            _Run(r, paragraph)
            for r in parse_xml(paragraph._element.xml.replace("fld", "r")).r_lst
        ]
    if len(runs) == 1:
        return runs[0]
    if len(runs) == 0:
        return None

    # Find the run with the most text
    run = max(runs, key=lambda x: len(x.text))
    run.text = paragraph.text

    # Remove other runs
    for r in runs:
        if r != run:
            r._r.getparent().remove(r._r)
    return run


def older_than(filepath: str, seconds: int = 10, wait: bool = False) -> bool:
    """
    Check if a file is older than a specified number of seconds.

    Args:
        filepath (str): The path to the file.
        seconds (int): The number of seconds to check against.
        wait (bool): Whether to wait for the file to exist.

    Returns:
        bool: True if the file is older than the specified number of seconds, False otherwise.
    """
    if not os.path.exists(filepath):
        while wait:
            logger.info("waiting for: %s", filepath)
            sleep(1)
            if os.path.exists(filepath):
                sleep(seconds)
                return True
        return False
    file_creation_time = os.path.getctime(filepath)
    current_time = time()
    return seconds < (current_time - file_creation_time)


def edit_distance(text1: str, text2: str) -> float:
    """
    Calculate the normalized edit distance between two strings.

    Args:
        text1 (str): The first string.
        text2 (str): The second string.

    Returns:
        float: The normalized edit distance (0.0 to 1.0, where 1.0 means identical).
    """
    if not text1 and not text2:
        return 1.0
    return 1 - Levenshtein.distance(text1, text2) / max(len(text1), len(text2))


def tenacity_log(retry_state: RetryCallState) -> None:
    """
    Log function for tenacity retries.

    Args:
        retry_state (RetryCallState): The retry state.
    """
    logger.warning("tenacity retry: %s", retry_state)
    traceback.print_tb(retry_state.outcome.exception().__traceback__)


def get_json_from_response(response: str) -> dict[str, Any]:
    """
    Extract JSON from a text response.

    Args:
        response (str): The response text.

    Returns:
        Dict[str, Any]: The extracted JSON.

    Raises:
        Exception: If JSON cannot be extracted from the response.
    """
    response = response.strip()

    try:
        return json.loads(response)
    except Exception:
        pass

    # Try to extract JSON from markdown code blocks
    l, r = response.rfind("```json"), response.rfind("```")
    if l != -1 and r != -1:
        json_obj = json_repair.loads(response[l + 7 : r].strip())
        if isinstance(json_obj, (dict, list)):
            return json_obj

    # Try to find JSON by looking for matching braces
    open_braces = []
    close_braces = []

    for i, char in enumerate(response):
        if char == "{" or char == "[":
            open_braces.append(i)
        elif char == "}" or char == "]":
            close_braces.append(i)

    for i, j in product(open_braces, reversed(close_braces)):
        if i > j:
            continue
        try:
            json_obj = json_repair.loads(response[i : j + 1])
            if isinstance(json_obj, (dict, list)):
                return json_obj
        except Exception:
            pass

    raise Exception("JSON not found in the given output", response)


# Create a tenacity decorator with custom settings
def tenacity_decorator(_func=None, *, wait: int = 3, stop: int = 5):
    def decorator(func):
        return retry(wait=wait_fixed(wait), stop=stop_after_attempt(stop))(func)

    if _func is None:
        # Called with arguments
        return decorator
    else:
        # Called without arguments
        return decorator(_func)


TABLE_CSS = """
table {
    border-collapse: collapse;  /* Merge borders */
    width: auto;               /* Width adapts to content */
    font-family: SimHei, Arial, sans-serif;  /* Font supporting Chinese characters */
    background: white;
}
th, td {
    border: 1px solid black;  /* Add borders */
    padding: 8px;             /* Cell padding */
    text-align: center;       /* Center text */
}
th {
    background-color: #f2f2f2; /* Header background color */
}
"""


# Convert Markdown to HTML
def get_chrome_path():
    """
    获取Chrome浏览器路径，兼容不同操作系统
    """
    if platform.system() == "Windows":
        possible_paths = [
            r"C:\Program Files\Google\Chrome\Application\chrome.exe",
            r"C:\Program Files (x86)\Google\Chrome\Application\chrome.exe",
            r"C:\Users\{}\AppData\Local\Google\Chrome\Application\chrome.exe".format(os.getenv("USERNAME", "")),
            # Edge浏览器作为备选
            r"C:\Program Files (x86)\Microsoft\Edge\Application\msedge.exe",
            r"C:\Program Files\Microsoft\Edge\Application\msedge.exe",
        ]

        for path in possible_paths:
            if os.path.exists(path):
                return path
        return None
    elif platform.system() == "Darwin":  # macOS
        possible_paths = [
            "/Applications/Google Chrome.app/Contents/MacOS/Google Chrome",
            "/Applications/Chromium.app/Contents/MacOS/Chromium",
        ]
        for path in possible_paths:
            if os.path.exists(path):
                return path
        return None
    else:  # Linux
        for cmd in ["google-chrome", "chromium-browser", "chromium"]:
            if which(cmd):
                return cmd
        return None

def markdown_table_to_image(markdown_text: str, output_path: str):
    """
    Convert a Markdown table to a cropped image, Windows compatible

    Args:
    markdown_text (str): Markdown text containing a table
    output_path (str): Output image path, defaults to 'table_cropped.png'

    Returns:
    str: The path of the generated image
    """
    html = markdown(markdown_text)
    assert "table" in html, "Failed to find table in markdown"

    parent_dir, basename = os.path.split(output_path)

    try:
        # 尝试获取Chrome路径
        chrome_path = get_chrome_path()

        if chrome_path:
            hti = Html2Image(
                disable_logging=True,
                output_path=parent_dir,
                custom_flags=["--no-sandbox", "--headless"],
                browser_executable=chrome_path
            )
        else:
            logger.warning("Chrome/Chromium not found, using default browser")
            hti = Html2Image(
                disable_logging=True,
                output_path=parent_dir,
                custom_flags=["--no-sandbox", "--headless"],
            )

        hti.browser.use_new_headless = None
        hti.screenshot(html_str=html, css_str=TABLE_CSS, save_as=basename)

        img = PILImage.open(output_path).convert("RGB")
        bbox = img.getbbox()
        assert (
            bbox is not None
        ), "Failed to capture the bbox, may be markdown table conversion failed"
        bbox = (0, 0, bbox[2] + 10, bbox[3] + 10)
        img.crop(bbox).save(output_path)
        return output_path

    except Exception as e:
        logger.error(f"Failed to create image from markdown table: {e}")
        # 创建一个简单的占位符图片
        img = PILImage.new('RGB', (400, 200), color='white')
        from PIL import ImageDraw, ImageFont
        draw = ImageDraw.Draw(img)
        try:
            font = ImageFont.load_default()
        except:
            font = None
        draw.text((10, 10), "Table conversion failed", fill='black', font=font)
        img.save(output_path)
        return output_path


@tenacity_decorator
def ppt_to_images(file: str, output_dir: str):
    """
    将PPT文件转换为图片，兼容Windows系统
    """
    if SOFFICE_CMD is None:
        raise RuntimeError("LibreOffice (soffice) is not available. Please install LibreOffice.")

    assert pexists(file), f"File {file} does not exist"
    if pexists(output_dir):
        logger.warning(f"ppt2images: {output_dir} already exists")
    os.makedirs(output_dir, exist_ok=True)

    with tempfile.TemporaryDirectory() as temp_dir:
        command_list = [
            SOFFICE_CMD,
            "--headless",
            "--convert-to",
            "pdf",
            file,
            "--outdir",
            temp_dir,
        ]

        # Windows下需要特殊处理路径
        if platform.system() == "Windows":
            # 确保路径使用正确的分隔符
            file = os.path.normpath(file)
            temp_dir = os.path.normpath(temp_dir)
            command_list = [
                SOFFICE_CMD,
                "--headless",
                "--convert-to",
                "pdf",
                file,
                "--outdir",
                temp_dir,
            ]

        logger.debug(f"Running command: {' '.join(command_list)}")
        process = subprocess.Popen(
            command_list,
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
            shell=(platform.system() == "Windows")  # Windows下使用shell
        )
        out, err = process.communicate()

        if process.returncode != 0:
            raise RuntimeError(f"LibreOffice conversion failed with error: {err.decode()}")

        for f in os.listdir(temp_dir):
            if not f.endswith(".pdf"):
                continue
            temp_pdf = pjoin(temp_dir, f)

            if not PDF2IMAGE_AVAILABLE:
                raise RuntimeError("pdf2image is not available. Please install poppler-utils.")

            try:
                # Windows下可能需要指定poppler路径
                if platform.system() == "Windows":
                    # 尝试常见的poppler路径
                    poppler_paths = [
                        r"C:\poppler\Library\bin",
                        r"C:\Program Files\poppler\bin",
                        None  # 使用系统PATH
                    ]

                    images = None
                    for poppler_path in poppler_paths:
                        try:
                            images = convert_from_path(temp_pdf, dpi=72, poppler_path=poppler_path)
                            break
                        except Exception as e:
                            if poppler_path is None:
                                raise e
                            continue
                else:
                    images = convert_from_path(temp_pdf, dpi=72)

                if images is None:
                    raise RuntimeError("Failed to convert PDF to images")

                for i, img in enumerate(images):
                    img.save(pjoin(output_dir, f"slide_{i+1:04d}.jpg"))
                return

            except Exception as e:
                logger.error(f"PDF to image conversion failed: {e}")
                if platform.system() == "Windows":
                    logger.error("Please ensure poppler is installed. See WINDOWS_SETUP.md for instructions.")
                raise

        raise RuntimeError(
            f"No PDF file was created in the temporary directory: {file}\n"
            f"Output: {out.decode()}\n"
            f"Error: {err.decode()}"
        )


@tenacity_decorator
async def ppt_to_images_async(file: str, output_dir: str):
    """
    异步将PPT文件转换为图片，兼容Windows系统
    """
    if SOFFICE_CMD is None:
        raise RuntimeError("LibreOffice (soffice) is not available. Please install LibreOffice.")

    assert pexists(file), f"File {file} does not exist"
    if pexists(output_dir):
        logger.debug(f"ppt2images: {output_dir} already exists")
    os.makedirs(output_dir, exist_ok=True)

    with tempfile.TemporaryDirectory() as temp_dir:
        command_list = [
            SOFFICE_CMD,
            "--headless",
            "--convert-to",
            "pdf",
            file,
            "--outdir",
            temp_dir,
        ]

        # Windows下需要特殊处理路径
        if platform.system() == "Windows":
            file = os.path.normpath(file)
            temp_dir = os.path.normpath(temp_dir)
            command_list = [
                SOFFICE_CMD,
                "--headless",
                "--convert-to",
                "pdf",
                file,
                "--outdir",
                temp_dir,
            ]

        logger.debug(f"Running async command: {' '.join(command_list)}")

        if platform.system() == "Windows":
            # Windows下使用shell模式
            process = await asyncio.create_subprocess_shell(
                " ".join(f'"{cmd}"' if " " in cmd else cmd for cmd in command_list),
                stdout=asyncio.subprocess.PIPE,
                stderr=asyncio.subprocess.PIPE,
            )
        else:
            # Linux/macOS使用exec模式
            process = await asyncio.create_subprocess_exec(
                *command_list,
                stdout=asyncio.subprocess.PIPE,
                stderr=asyncio.subprocess.PIPE,
            )

        stdout, stderr = await process.communicate()
        if process.returncode != 0:
            raise RuntimeError(f"LibreOffice conversion failed with error: {stderr.decode()}")

        for f in os.listdir(temp_dir):
            if not f.endswith(".pdf"):
                continue
            temp_pdf = pjoin(temp_dir, f)

            if not PDF2IMAGE_AVAILABLE:
                raise RuntimeError("pdf2image is not available. Please install poppler-utils.")

            try:
                # Windows下可能需要指定poppler路径
                if platform.system() == "Windows":
                    # 尝试常见的poppler路径
                    poppler_paths = [
                        r"C:\path\poppler-24.08.0\Library\bin",  # 用户配置的路径
                        r"C:\poppler\Library\bin",
                        r"C:\Program Files\poppler\bin",
                        None  # 使用系统PATH
                    ]

                    images = None
                    for poppler_path in poppler_paths:
                        try:
                            images = convert_from_path(temp_pdf, dpi=72, poppler_path=poppler_path)
                            break
                        except Exception as e:
                            if poppler_path is None:
                                raise e
                            continue
                else:
                    images = convert_from_path(temp_pdf, dpi=72)

                if images is None:
                    raise RuntimeError("Failed to convert PDF to images")

                for i, img in enumerate(images):
                    img.save(pjoin(output_dir, f"slide_{i+1:04d}.jpg"))
                return

            except Exception as e:
                logger.error(f"PDF to image conversion failed: {e}")
                if platform.system() == "Windows":
                    logger.error("Please ensure poppler is installed. See WINDOWS_SETUP.md for instructions.")
                raise

        raise RuntimeError(
            f"No PDF file was created in the temporary directory: {file}\n"
            f"Output: {stdout.decode()}\n"
            f"Error: {stderr.decode()}"
        )


def parsing_image(image: Image, image_path: str) -> str:
    # Handle WMF images (PDFs)
    if image.ext == "wmf":
        image_path = image_path.replace(".wmf", ".jpg")
        if not pexists(image_path):
            wmf_to_images(image.blob, image_path)
    # Check for supported image types
    elif image.ext not in IMAGE_EXTENSIONS:
        raise ValueError(f"Unsupported image type {image.ext}")

    # Save image if it doesn't exist
    if not pexists(image_path):
        with open(image_path, "wb") as f:
            f.write(image.blob)
    return image_path


@tenacity_decorator
def wmf_to_images(blob: bytes, filepath: str):
    """
    将WMF格式图片转换为JPG，兼容Windows系统
    """
    if SOFFICE_CMD is None:
        raise RuntimeError("LibreOffice (soffice) is not available. Please install LibreOffice.")

    if not filepath.endswith(".jpg"):
        raise ValueError("filepath must end with .jpg")
    dirname = os.path.dirname(filepath)
    basename = os.path.basename(filepath).removesuffix(".jpg")

    with tempfile.TemporaryDirectory() as temp_dir:
        wmf_path = pjoin(temp_dir, f"{basename}.wmf")
        with open(wmf_path, "wb") as f:
            f.write(blob)

        command_list = [
            SOFFICE_CMD,
            "--headless",
            "--convert-to",
            "jpg",
            wmf_path,
            "--outdir",
            dirname,
        ]

        # Windows下需要特殊处理
        if platform.system() == "Windows":
            wmf_path = os.path.normpath(wmf_path)
            dirname = os.path.normpath(dirname)
            command_list = [
                SOFFICE_CMD,
                "--headless",
                "--convert-to",
                "jpg",
                wmf_path,
                "--outdir",
                dirname,
            ]

        logger.debug(f"Converting WMF: {' '.join(command_list)}")
        subprocess.run(
            command_list,
            check=True,
            stdout=subprocess.DEVNULL,
            shell=(platform.system() == "Windows")
        )

    assert pexists(filepath), f"File {filepath} does not exist"


def parse_groupshape(groupshape: GroupShape) -> list[dict[str, Length]]:
    """
    Parse a group shape to get the bounds of its child shapes.

    Args:
        groupshape (GroupShape): The group shape to parse.

    Returns:
        List[Dict[str, Length]]: The bounds of the child shapes.

    Raises:
        AssertionError: If the input is not a GroupShape.
    """
    assert isinstance(groupshape, GroupShape), "Input must be a GroupShape"

    # Get group bounds
    group_top_left_x = groupshape.left
    group_top_left_y = groupshape.top
    group_width = groupshape.width
    group_height = groupshape.height

    # Get shape bounds
    shape_top_left_x = min([sp.left for sp in groupshape.shapes])
    shape_top_left_y = min([sp.top for sp in groupshape.shapes])
    shape_width = (
        max([sp.left + sp.width for sp in groupshape.shapes]) - shape_top_left_x
    )
    shape_height = (
        max([sp.top + sp.height for sp in groupshape.shapes]) - shape_top_left_y
    )

    # Calculate bounds for each shape in the group
    group_shape_xy = []
    for sp in groupshape.shapes:
        group_shape_left = (
            sp.left - shape_top_left_x
        ) * group_width / shape_width + group_top_left_x
        group_shape_top = (
            sp.top - shape_top_left_y
        ) * group_height / shape_height + group_top_left_y
        group_shape_width = sp.width * group_width / shape_width
        group_shape_height = sp.height * group_height / shape_height

        group_shape_xy.append(
            {
                "left": Length(group_shape_left),
                "top": Length(group_shape_top),
                "width": Length(group_shape_width),
                "height": Length(group_shape_height),
            }
        )

    return group_shape_xy


def is_primitive(obj: Any) -> bool:
    """
    Check if an object is a primitive type or a collection of primitive types.

    Args:
        obj (Any): The object to check.

    Returns:
        bool: True if the object is a primitive type or a collection of primitive types, False otherwise.
    """
    if isinstance(obj, (list, tuple, set, frozenset)):
        return all(is_primitive(item) for item in obj)

    return isinstance(
        obj, (int, float, complex, bool, str, bytes, bytearray, type(None))
    )


DEFAULT_EXCLUDE: set[str] = {"element", "language_id", "ln", "placeholder_format"}


def dict_to_object(
    dict_obj: dict[str, Any], obj: Any, exclude: Optional[set[str]] = None
) -> None:
    """
    Apply dictionary values to an object.

    Args:
        dict_obj (Dict[str, Any]): The dictionary with values to apply.
        obj (Any): The object to apply values to.
        exclude (Optional[Set[str]]): The keys to exclude.
    """
    if exclude is None:
        exclude = set()

    for key, value in dict_obj.items():
        if key not in exclude and value is not None:
            setattr(obj, key, value)


def package_join(*paths: str) -> str:
    """
    Join paths with the appropriate separator for the platform.

    Args:
        *paths: The paths to join.

    Returns:
        str: The joined path.
    """
    _dir = pdirname(__file__)
    return pjoin(_dir, *paths)


class Config:
    """
    Configuration class for the application.
    """

    def __init__(
        self,
        rundir: Optional[str] = None,
        session_id: Optional[str] = None,
    ):
        """
        Initialize the configuration.

        Args:
            rundir (Optional[str]): The run directory.
            session_id (Optional[str]): The session ID.
            debug (bool): Whether to enable debug mode.
        """
        if rundir is not None:
            self.set_rundir(rundir)
        elif session_id is not None:
            self.set_session(session_id)
        else:
            raise ValueError("No session ID or run directory provided")

    def set_session(self, session_id: str) -> None:
        """
        Set the session ID and update the run directory.

        Args:
            session_id (str): The session ID.
        """
        self.session_id = session_id
        self.set_rundir(f"./runs/{session_id}")

    def set_rundir(self, rundir: str) -> None:
        """
        Set the run directory and create necessary subdirectories.

        Args:
            rundir (str): The run directory.
        """
        # 确保路径使用正确的分隔符
        self.RUN_DIR = os.path.normpath(rundir)
        self.IMAGE_DIR = pjoin(self.RUN_DIR, "images")

        for the_dir in [self.RUN_DIR, self.IMAGE_DIR]:
            os.makedirs(the_dir, exist_ok=True)

    def set_debug(self, debug: bool) -> None:
        """
        Set the debug mode.

        Args:
            debug (bool): Whether to enable debug mode.
        """
        self.DEBUG = debug

    def remove_rundir(self) -> None:
        """
        Remove the run directory and its subdirectories.
        """
        if pexists(self.RUN_DIR):
            shutil.rmtree(self.RUN_DIR)
        if pexists(self.IMAGE_DIR):
            shutil.rmtree(self.IMAGE_DIR)

    def __repr__(self) -> str:
        """
        Get a string representation of the configuration.

        Returns:
            str: A string representation of the configuration.
        """
        attrs = []
        for attr in dir(self):
            if not attr.startswith("_") and not callable(getattr(self, attr)):
                attrs.append(f"{attr}={getattr(self, attr)}")
        return f"Config({', '.join(attrs)})"


# Path utility functions
pjoin = os.path.join
pexists = os.path.exists
pbasename = os.path.basename
pdirname = os.path.dirname
